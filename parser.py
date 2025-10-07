# app/services/parser.py

import os
import pdfplumber
import fitz
import docx
from pptx import Presentation
import pandas as pd
from pdf2image import convert_from_path
from PIL import Image
import easyocr
import numpy as np 

# Initialize EasyOCR once (globally)
try:
    ocr_reader = easyocr.Reader(['en'], gpu=False)
except Exception as e:
    print(f"Warning: EasyOCR initialization failed: {e}")
    ocr_reader = None


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF. Use OCR if scanned PDF."""
    text = ""

    # Try pdfplumber (text-based PDF)
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed: {e}")

    # Try PyMuPDF (fitz) if pdfplumber fails
    if not text.strip():
        try:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text()
            doc.close()
        except Exception as e:
            print(f"PyMuPDF failed: {e}")

    # OCR fallback for scanned PDFs (Requires Poppler installed for pdf2image)
    if not text.strip() and ocr_reader:
        try:
            pages = convert_from_path(file_path)
            for page in pages:
                # Use EasyOCR to extract text from each page image
                ocr_result = ocr_reader.readtext(np.array(page), detail=0)
                text += "\n".join(ocr_result)
        except Exception as e:
            print(f"OCR failed for PDF: {e}")

    return text


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX files"""
    try:
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        print(f"DOCX extraction failed: {e}")
        return ""


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT files"""
    try:
        with open(file_path, "r", encoding="utf-8", errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"TXT extraction failed: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX files"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"PPTX extraction failed: {e}")
        return ""


def extract_text_from_excel(file_path: str) -> str:
    """Extract text from Excel/CSV files"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".xlsx":
            df = pd.read_excel(file_path)
        elif ext == ".csv":
            df = pd.read_csv(file_path)
        else:
            raise ValueError(f"Unsupported excel extension: {ext}")
        # Convert all columns to string for LLM input
        return df.to_string()
    except Exception as e:
        print(f"Excel extraction failed: {e}")
        return ""


def extract_text_from_image(file_path: str) -> str:
    """Extract text from image files using OCR"""
    if not ocr_reader:
        return "ERROR: OCR not available"
    
    try:
        # Use EasyOCR to extract text from the image
        text = "\n".join(ocr_reader.readtext(file_path, detail=0))
        return text
    except Exception as e:
        print(f"Image OCR failed: {e}")
        return ""


def extract_text_from_file(file_path: str) -> str:
    """Detect file type and extract text accordingly."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        return extract_text_from_docx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(file_path)
    elif ext in [".csv", ".xlsx"]:
        return extract_text_from_excel(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        return extract_text_from_image(file_path)
    else:
        # Fallback for unexpected file types
        raise ValueError(f"ERROR: Unsupported file type: {ext}")